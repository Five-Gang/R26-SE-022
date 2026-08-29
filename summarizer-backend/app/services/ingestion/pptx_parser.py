from __future__ import annotations
from typing import Optional, Union
"""PPTX Parser Service — extracts content from PowerPoint files.

Preserves slide ordering as pedagogical sequence and extracts
text from shapes, tables, and speaker notes.
"""

from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


@dataclass
class ExtractedSlide:
    """Content extracted from a single slide."""

    slide_number: int
    title: str = ""
    body_text: str = ""
    speaker_notes: str = ""
    has_images: bool = False
    has_tables: bool = False
    has_charts: bool = False
    table_data: list[list[str]] = field(default_factory=list)
    image_count: int = 0

    @property
    def full_text(self) -> str:
        """Combined text from title, body, and notes."""
        parts = []
        if self.title:
            parts.append(f"## {self.title}")
        if self.body_text:
            parts.append(self.body_text)
        if self.table_data:
            parts.append(self._format_tables())
        if self.speaker_notes:
            parts.append(f"\n[Speaker Notes]: {self.speaker_notes}")
        return "\n\n".join(parts)

    def _format_tables(self) -> str:
        """Format extracted table data as markdown."""
        if not self.table_data:
            return ""
        lines = []
        for i, row in enumerate(self.table_data):
            lines.append("| " + " | ".join(row) + " |")
            if i == 0:
                lines.append("| " + " | ".join("---" for _ in row) + " |")
        return "\n".join(lines)

    @property
    def char_count(self) -> int:
        return len(self.full_text)


@dataclass
class PPTXExtractionResult:
    """Complete extraction result from a PPTX file."""

    slides: list[ExtractedSlide] = field(default_factory=list)
    full_text: str = ""
    slide_count: int = 0
    metadata: dict = field(default_factory=dict)

    @property
    def total_chars(self) -> int:
        return sum(s.char_count for s in self.slides)


class PPTXParser:
    """Extracts structured content from PowerPoint PPTX files.

    Preserves:
    - Slide ordering (pedagogical sequence)
    - Slide titles (section markers)
    - Body text (from text boxes and shapes)
    - Speaker notes (often contain detailed explanations)
    - Table data (formatted as markdown)
    - Image detection flags
    """

    def extract(self, file_path: Union[str, Path]) -> PPTXExtractionResult:
        """Extract content from a PPTX file.

        Args:
            file_path: Path to the PPTX file.

        Returns:
            PPTXExtractionResult with extracted slides and metadata.
        """
        prs = Presentation(str(file_path))
        result = PPTXExtractionResult()
        result.slide_count = len(prs.slides)
        result.metadata = self._extract_metadata(prs)

        for slide_num, slide in enumerate(prs.slides, start=1):
            extracted = self._extract_slide(slide, slide_num)
            result.slides.append(extracted)

        # Combine full text with slide separators
        result.full_text = "\n\n---\n\n".join(
            f"[Slide {s.slide_number}]\n{s.full_text}"
            for s in result.slides
            if s.full_text.strip()
        )

        return result

    def extract_from_bytes(self, content: bytes) -> PPTXExtractionResult:
        """Extract content from PPTX bytes (for in-memory processing)."""
        import io

        prs = Presentation(io.BytesIO(content))
        result = PPTXExtractionResult()
        result.slide_count = len(prs.slides)
        result.metadata = self._extract_metadata(prs)

        for slide_num, slide in enumerate(prs.slides, start=1):
            extracted = self._extract_slide(slide, slide_num)
            result.slides.append(extracted)

        result.full_text = "\n\n---\n\n".join(
            f"[Slide {s.slide_number}]\n{s.full_text}"
            for s in result.slides
            if s.full_text.strip()
        )

        return result

    def _extract_slide(self, slide, slide_number: int) -> ExtractedSlide:
        """Extract content from a single slide, robustly handling all shape types."""
        extracted = ExtractedSlide(slide_number=slide_number)

        title_text = []
        body_texts = []
        image_count = 0

        # First, try to get the title cleanly from the slide title shape
        try:
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                t = slide.shapes.title.text_frame.text.strip()
                if t:
                    title_text.append(t)
        except Exception:
            pass

        title_shape_id = None
        try:
            if slide.shapes.title:
                title_shape_id = slide.shapes.title.shape_id
        except Exception:
            pass

        for shape in slide.shapes:
            try:
                # Skip the title shape — already handled above
                if title_shape_id is not None and shape.shape_id == title_shape_id:
                    continue

                # Text frames
                if shape.has_text_frame:
                    text = self._extract_text_frame(shape.text_frame)
                    if text.strip():
                        body_texts.append(text)

                # Tables
                if shape.has_table:
                    extracted.has_tables = True
                    table_rows = self._extract_table(shape.table)
                    extracted.table_data.extend(table_rows)

                # Images
                try:
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        image_count += 1
                        extracted.has_images = True
                except Exception:
                    pass

                # Charts
                try:
                    if shape.has_chart:
                        extracted.has_charts = True
                except Exception:
                    pass

            except Exception:
                # Never crash on a single shape — skip it and continue
                continue

        extracted.title = " ".join(title_text).strip()
        extracted.body_text = "\n".join(body_texts).strip()
        extracted.image_count = image_count

        # Speaker notes
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes and notes.lower() != "click to add notes":
                    extracted.speaker_notes = notes
        except Exception:
            pass

        return extracted

    def _extract_text_frame(self, text_frame) -> str:
        """Extract text from a text frame, preserving bullet structure."""
        lines = []
        for paragraph in text_frame.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            # Add bullet prefix based on indent level
            indent = paragraph.level if paragraph.level else 0
            prefix = "  " * indent + "- " if indent > 0 else ""
            lines.append(f"{prefix}{text}")
        return "\n".join(lines)

    def _extract_table(self, table) -> list[list[str]]:
        """Extract table data as a list of rows."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)
        return rows

    def _extract_metadata(self, prs: Presentation) -> dict:
        """Extract PPTX presentation metadata."""
        props = prs.core_properties
        return {
            "title": props.title or "",
            "author": props.author or "",
            "subject": props.subject or "",
            "keywords": props.keywords or "",
            "last_modified_by": props.last_modified_by or "",
            "created": str(props.created) if props.created else "",
            "modified": str(props.modified) if props.modified else "",
        }
