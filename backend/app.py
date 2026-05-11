#!/usr/bin/env python3
"""
Component 3: Lecture Summarizer & Flashcard Generator (ML-Enhanced)
Uses trained RandomForest model for importance scoring + module-specific adaptation

Models:
- universal_model.pkl: Trained RandomForest for importance classification
- vectorizer.pkl: TF-IDF vectorizer for feature extraction
- module_config.json: Module-specific learning objectives and keywords

Run: pip install flask flask-cors scikit-learn nltk PyPDF2 python-pptx numpy anthropic
"""

from flask import Flask, request, jsonify
from flask_cors import CORS
import os
import re
import json
import io
import pickle
import sys
from pathlib import Path

import numpy as np
from nltk.tokenize import sent_tokenize
import nltk

# Add scripts folder to path for importing module_score_booster
sys.path.insert(0, str(Path(__file__).parent.parent / 'scripts'))

app = Flask(__name__)
CORS(app)

# ─── NLP & File Format Imports ─────────────────────────────────────────────────
try:
    import nltk
    for pkg in ['punkt', 'stopwords', 'punkt_tab']:
        try:
            nltk.data.find(f'tokenizers/{pkg}')
        except LookupError:
            nltk.download(pkg, quiet=True)
    NLP_AVAILABLE = True
except ImportError:
    NLP_AVAILABLE = False

try:
    from PyPDF2 import PdfReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.ensemble import RandomForestClassifier
    from sklearn.metrics.pairwise import cosine_similarity
    ML_AVAILABLE = True
except ImportError:
    ML_AVAILABLE = False

# ─── Global Models & Config ───────────────────────────────────────────────────
MODEL_PATH = Path(__file__).parent / 'universal_model.pkl'
VECTORIZER_PATH = Path(__file__).parent / 'vectorizer.pkl'
MODULE_CONFIG_PATH = Path(__file__).parent / 'module_config.json'
METADATA_PATH = Path(__file__).parent / 'universal_metadata.json'

MODEL = None
VECTORIZER = None
MODULE_CONFIG = None
OPTIMAL_THRESHOLD = 0.5   # default; overridden by metadata if present

def load_models():
    """Load trained model, vectorizer, and module config at startup"""
    global MODEL, VECTORIZER, MODULE_CONFIG
    
    print("Loading models...")
    
    # Try to load trained model
    if MODEL_PATH.exists():
        try:
            with open(MODEL_PATH, 'rb') as f:
                MODEL = pickle.load(f)
            print(f"✅ Loaded model: {MODEL_PATH}")
        except Exception as e:
            print(f"❌ Error loading model: {e}")
    
    # Try to load vectorizer
    if VECTORIZER_PATH.exists():
        try:
            with open(VECTORIZER_PATH, 'rb') as f:
                VECTORIZER = pickle.load(f)
            print(f"✅ Loaded vectorizer: {VECTORIZER_PATH}")
        except Exception as e:
            print(f"❌ Error loading vectorizer: {e}")
    
    # Load module config
    if MODULE_CONFIG_PATH.exists():
        try:
            with open(MODULE_CONFIG_PATH, 'r') as f:
                MODULE_CONFIG = json.load(f)
            print(f"✅ Loaded module config: {len(MODULE_CONFIG.get('modules', []))} modules")
        except Exception as e:
            print(f"❌ Error loading module config: {e}")
    
    # Load optimal threshold from metadata
    global OPTIMAL_THRESHOLD
    if METADATA_PATH.exists():
        try:
            with open(METADATA_PATH, 'r') as f:
                meta = json.load(f)
            OPTIMAL_THRESHOLD = meta.get('optimal_threshold', 0.5)
            print(f'✅ Optimal threshold loaded: {OPTIMAL_THRESHOLD:.3f}')
        except Exception as e:
            print(f'⚠️  Could not load metadata: {e}')

    if not MODEL or not VECTORIZER:
        print("\n⚠️  Trained model not found. Using fallback TextRank mode.")
        print("To train model, run:")
        print("  python3 scripts/extract_sentences_multi_module.py")
        print("  python3 scripts/auto_label_multi_module.py")
        print("  python3 scripts/train_universal_model.py")


# Load models at startup
load_models()

# ─── Import Score Booster ─────────────────────────────────────────────────────
try:
    from module_score_booster import (
        boost_scores_for_module,
        get_top_sentences,
        get_module_statistics,
        map_scores_to_learning_objectives
    )
    BOOSTER_AVAILABLE = True
except ImportError:
    BOOSTER_AVAILABLE = False
    print("Warning: module_score_booster not available")


# ─── Feature Extraction for ML Model ───────────────────────────────────────────
# Must stay byte-for-byte identical to extract_structural_features() in the
# Colab training notebook (cell-14).  Any change here must also be made there.

_FUNC_TECH = re.compile(
    r'\b(data|system|network|protocol|server|client|query|table|process|thread|'
    r'memory|file|algorithm|function|class|object|database|schema|index|'
    r'transaction|security|architecture|layer|node|packet|buffer|cache|'
    r'stack|queue|tree|graph|cpu|os|disk|page|segment|register|bit|byte|'
    r'compiler|interpreter|kernel|scheduler|mutex|semaphore|deadlock|lock|'
    r'primary key|foreign key|normalization|acid|atomicity|consistency|isolation|durability|'
    r'concurrency|replication|partition|encryption|authentication|authorization|'
    r'bandwidth|latency|throughput|subnet|routing|switching|firewall|'
    r'heap|pointer|recursion|polymorphism|inheritance|encapsulation|abstraction|'
    r'sorting|hashing|linked list|binary tree|relational|tuple|attribute|domain|'
    r'constraint|trigger|view|stored procedure|join|aggregate|subquery)\b',
    re.IGNORECASE
)
_PRONOUN_START = re.compile(r'^(he|she|it|they|we|you|i)\s', re.IGNORECASE)


def _is_definition(s):
    """Tightened definition check — mirrors label_sentence rule 11."""
    if _PRONOUN_START.match(s):
        return False
    if re.match(
        r'^[A-Z][a-zA-Z\s\-/()\d]{0,45}?\b'
        r'(is defined as|is known as|refers to|can be defined as|means|is called)\b', s
    ):
        return True
    if re.match(
        r'^([A-Z][a-z]+(?:\s+[A-Za-z]+){0,3})\s+(is|are)\s+(a|an|the|used|considered|defined|known|responsible)\b', s
    ):
        return True
    if re.match(
        r'^(A|An|The)\s+[a-z][a-z\-]+(?:\s+[a-z]+){0,2}\s+(is|are)\s+(a|an|the|used|considered|defined)\b', s
    ):
        return True
    return False


def extract_structural_features(sentence, position_ratio=0.5):
    """
    21 structural features — must stay in sync with the Colab training notebook.
    """
    s = sentence
    sl = sentence.lower()
    words = s.split()
    word_count = len(words)

    features = [
        # 0  Basic text stats
        word_count,
        np.mean([len(w) for w in words]) if words else 0,
        sum(1 for c in s if c.isupper()) / max(len(s), 1),       # capital ratio
        sum(1 for c in s if c.isdigit()) / max(len(s), 1),       # digit ratio
        sum(1 for c in s if c in '.,;:!?') / max(len(s), 1),     # punct ratio

        # 5  Definition (tightened — pronoun-start excluded)
        1 if _is_definition(s) else 0,

        # 6  Acronym expansion e.g. "SQL (Structured …"
        1 if re.search(r'[A-Z]{2,}\s*\([A-Z][a-z]', s) else 0,

        # 7  Composition (includes? added — "ACID includes Atomicity…")
        1 if (re.search(r'\b(consists? of|is composed of|comprises?|includes?|contain)\b', sl)
              and (_FUNC_TECH.search(sl) or re.search(r'\b[A-Z]{2,}\b', s))) else 0,

        # 8  Enumeration
        1 if re.search(r'there are (\d+|two|three|four|five|six|several)', sl) else 0,

        # 9  Conclusion
        1 if re.search(r'\b(therefore|thus|hence|in summary|in conclusion)\b', sl) else 0,

        # 10 Advantages / disadvantages
        1 if re.search(r'\b(advantage|disadvantage|benefit|drawback|limitation)\b', sl) else 0,

        # 11 Causal connector
        1 if re.search(r'\b(because|since|due to|in order to|so that)\b', sl) else 0,

        # 12 Example marker
        1 if re.search(r'\b(for example|e\.g\.|such as|for instance)\b', sl) else 0,

        # 13 Comparison
        1 if re.search(r'\b(difference between|compared to|in contrast|unlike)\b', sl) else 0,

        # 14 Functional description — TIGHTENED: also requires a technical noun
        1 if (re.search(r'\b(is used (to|for)|allows|enables|provides|ensures)\b', sl)
              and _FUNC_TECH.search(sl)) else 0,

        # 15 Acronym count
        len(re.findall(r'\b[A-Z]{2,}\b', s)),

        # 16 Technical term count (same vocab as notebook _FUNC_TECH)
        len(_FUNC_TECH.findall(sl)),

        # 17 Position in document
        position_ratio,

        # 18 Has email
        1 if re.search(r'@\w+\.\w+', s) else 0,

        # 19 Boilerplate
        1 if re.search(r'copyright|all rights reserved', sl) else 0,

        # 20 Reference line
        1 if re.match(r'^(see|refer to|read|check)\s', sl) else 0,
    ]
    return features


def score_sentences_with_model(sentences):
    """Score sentences using trained model"""
    if MODEL is None or VECTORIZER is None:
        return np.ones(len(sentences)) * 0.5

    try:
        tfidf_features = VECTORIZER.transform(sentences).toarray()
        structural = np.array([
            extract_structural_features(s, i / max(len(sentences) - 1, 1))
            for i, s in enumerate(sentences)
        ])
        features = np.hstack([tfidf_features, structural])

        # Use optimal threshold from training to produce a calibrated 0/1 signal,
        # then return raw probabilities for downstream blending
        raw_probs = MODEL.predict_proba(features)[:, 1]
        # Rescale so that optimal_threshold maps to 0.5 (preserves blend math)
        scores = raw_probs / (2.0 * OPTIMAL_THRESHOLD + 1e-9) * 1.0
        scores = np.clip(scores, 0.0, 1.0)
        return scores
    except Exception as e:
        print(f"Error scoring sentences: {e}")
        return np.ones(len(sentences)) * 0.5


# ─── Text Extraction ──────────────────────────────────────────────────────────
def extract_text_from_pdf(file_bytes):
    if not PDF_AVAILABLE:
        return ""
    try:
        reader = PdfReader(io.BytesIO(file_bytes))
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + " "
        return text
    except Exception as e:
        print(f"Error extracting PDF: {e}")
        return ""


def extract_text_from_pptx(file_bytes):
    if not PPTX_AVAILABLE:
        return ""
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    texts.append(shape.text)
        return ' '.join(texts)
    except Exception as e:
        print(f"Error extracting PPTX: {e}")
        return ""


def extract_sentences_from_text(text):
    """Extract and clean sentences from text, filtering PDF noise."""
    try:
        text = re.sub(r'\x0c', '\n\n', text)          # form-feed → paragraph break
        text = re.sub(r'[\x00-\x08\x0b\x0e-\x1f]', '', text)
        text = re.sub(r'/g\d+', '', text)
        text = re.sub(r'\(cid:\d+\)', '', text)

        # ── Strip module-code header lines BEFORE tokenising ──────────────────
        # Catches: "IT2060 - Operating Systems ... Lecture 02 ..."
        text = re.sub(r'[A-Z]{2}\d{3,4}\s*[-–]\s*[^\n]*', ' ', text)
        # Catches: "IT2060 | ... | Lecture 02 | ..."
        text = re.sub(r'[A-Z]{2}\d{3,4}\s*\|[^\n]*', ' ', text)
        # Catches standalone "Lecture 02 Introduction to ..." header lines
        text = re.sub(r'\bLecture\s+\d+[^\n.!?]{0,80}\n', ' ', text)

        text = re.sub(r'\n{2,}', '. ', text)
        text = re.sub(r'(?<=[a-z])\n(?=[a-z])', ' ', text)
        text = re.sub(r'\s+', ' ', text).strip()

        sentences = sent_tokenize(text)
        cleaned = []
        for s in sentences:
            s = s.strip()

            # ── Strip residual header fragments ───────────────────────────────
            s = re.sub(r'^[\s•\-\*▪▸►‣◦]+', '', s).strip()
            s = re.sub(r'Module\s+Code\s*\|.*$', '', s, flags=re.IGNORECASE).strip()
            s = re.sub(r'Lecture\s+Title\s*\|.*$', '', s, flags=re.IGNORECASE).strip()
            s = re.sub(r'\|\s*Lecturer\s*(Name)?\s*\.?$', '', s, flags=re.IGNORECASE).strip()
            s = re.sub(r'[A-Z]{2}\d{3,4}\s*\|[^.!?]*', '', s).strip()
            s = re.sub(r'[A-Z]{2}\d{3,4}\s*[-–][^.!?]{0,60}', '', s).strip()
            s = re.sub(r'\bLecture\s+\d+\b[^.!?]{0,40}', '', s).strip()
            s = re.sub(r'Shashika\s*Lokuliyana\s*\d*', '', s, flags=re.IGNORECASE).strip()
            s = re.sub(r'Faculty\s*of\s*Computing[^.]*', '', s).strip()
            s = re.sub(r'^\d+\s+', '', s).strip()
            s = re.sub(r'\s+\d+\s*$', '', s).strip()

            # Instructor name filter
            if any(n in s.lower() for n in ['sanvitha', 'kasthuriarachchi']):
                continue

            words = s.split()
            wc = len(words)

            # ── Length gate — match training data range (6–80 words) ──────────
            if wc < 6 or wc > 80:
                continue

            # ── Alpha ratio ───────────────────────────────────────────────────
            alpha_ratio = sum(1 for c in s if c.isalpha()) / max(len(s), 1)
            if alpha_ratio < 0.55:
                continue

            if s.count('...') > 2 or s.count('_') > 5:
                continue

            # ── Slide table-of-contents / outline detector ────────────────────
            # These are slide outline pages where bullet headings got concatenated:
            # "Storage Structure Booting Up Process Multiprocessor System Interrupts..."
            # Signal: very high ratio of Title-Case words with no sentence structure.
            title_case_words = [w for w in words if w and w[0].isupper() and w.isalpha()]
            title_case_ratio = len(title_case_words) / max(wc, 1)
            # More than 55% Title-Case AND sentence is 15+ words → slide outline
            if title_case_ratio > 0.55 and wc >= 15:
                # Allow it only if it has a clear sentence verb
                has_sent_verb = bool(re.search(
                    r'\b(is|are|was|were|has|have|can|will|must|should|defined|used|'
                    r'provides?|requires?|ensures?|consists?|includes?|refers|means)\b',
                    s, re.IGNORECASE
                ))
                if not has_sent_verb:
                    continue

            # ── All-caps word ratio (diagram / flowchart labels) ──────────────
            all_caps = [w for w in words if w.isupper() and len(w) >= 2]
            if len(all_caps) / max(wc, 1) > 0.4:
                continue

            # ── Jammed words from bad PDF extraction ──────────────────────────
            long_words = [w for w in words if len(w) > 20 and not w.startswith('http')]
            if len(long_words) > 2:
                continue

            # ── Must have at least one verb to be a real sentence ─────────────
            has_verb = bool(re.search(
                r'\b(is|are|was|were|has|have|can|will|must|should|provides?|requires?|'
                r'allows?|ensures?|creates?|stores?|manages?|consists?|includes?|defines?|'
                r'uses?|enables?|contains?|represents?|implements?|occurs?|refers?)\b',
                s, re.IGNORECASE
            ))
            if not has_verb:
                continue

            cleaned.append(s)
        return cleaned
    except Exception as e:
        print(f"Error tokenizing: {e}")
        return [s.strip() for s in text.split('. ') if len(s.split()) >= 5][:50]


# ─── Flashcard Generation ─────────────────────────────────────────────────────
def generate_flashcards_from_sentences(sentences, module_id, learning_objectives):
    """Generate flashcards from selected sentences with varied question types."""
    flashcards = []

    definition_patterns = [
        (r'(.+?)\s+(?:is|are|refers to|means|defined as)\s+(.+)', 'definition'),
        (r'(.+?)\s+(?:consists? of|includes?|contains?)\s+(.+)', 'composition'),
    ]

    for idx, sentence in enumerate(sentences[:15]):
        sentence = sentence.strip()
        words = sentence.split()

        if len(words) < 5:
            continue

        sentence_lower = sentence.lower()
        question = None
        difficulty = 'medium'

        for pattern, ptype in definition_patterns:
            match = re.match(pattern, sentence, re.IGNORECASE)
            if match:
                subject = match.group(1).strip()
                if len(subject.split()) <= 6:
                    if ptype == 'definition':
                        question = f"What is {subject}?"
                    else:
                        question = f"What does {subject} consist of?"
                    difficulty = 'easy'
                break

        if not question:
            if any(kw in sentence_lower for kw in ['step', 'first', 'then', 'process', 'procedure']):
                question = f"Describe the process: {' '.join(words[:6])}..."
                difficulty = 'medium'
            elif any(kw in sentence_lower for kw in ['advantage', 'benefit', 'important', 'crucial']):
                question = f"Why is this important: {' '.join(words[:5])}...?"
                difficulty = 'medium'
            elif any(kw in sentence_lower for kw in ['difference', 'compare', 'versus', 'unlike']):
                question = f"Compare and contrast: {' '.join(words[:6])}..."
                difficulty = 'hard'
            elif any(kw in sentence_lower for kw in ['example', 'such as', 'e.g', 'instance']):
                question = f"Give an example related to: {' '.join(words[:5])}..."
                difficulty = 'easy'
            elif any(kw in sentence_lower for kw in ['formula', 'equation', 'calculate']):
                question = f"What is the formula/method for: {' '.join(words[:5])}...?"
                difficulty = 'hard'
            else:
                important_words = [w for w in words if len(w) > 4 and w[0].isupper()]
                if important_words:
                    question = f"Explain: {' '.join(important_words[:3])}"
                else:
                    question = f"What is the key idea: {' '.join(words[:6])}...?"
                difficulty = 'medium'

        lo_index = idx % len(learning_objectives) if learning_objectives else 0
        lo_id = f"LO{lo_index + 1}"
        lo_text = learning_objectives[lo_index] if learning_objectives and lo_index < len(learning_objectives) else f"Learning Objective {lo_index + 1}"

        flashcards.append({
            'id': len(flashcards) + 1,
            'question': question,
            'answer': sentence,
            'module': module_id,
            'learning_objective': lo_id,
            'lo_text': lo_text,
            'difficulty': difficulty
        })

    return flashcards


# ─── API Routes ───────────────────────────────────────────────────────────────
@app.route('/api/health', methods=['GET'])
def health():
    """Check system status"""
    return jsonify({
        'status': 'ok',
        'nlp_available': NLP_AVAILABLE,
        'pdf_available': PDF_AVAILABLE,
        'pptx_available': PPTX_AVAILABLE,
        'ml_model_available': MODEL is not None,
        'module_config_available': MODULE_CONFIG is not None
    })


@app.route('/api/modules', methods=['GET'])
def get_modules():
    """Get list of available modules with learning objectives"""
    if MODULE_CONFIG is None:
        return jsonify({'error': 'Module configuration not loaded'}), 500
    
    modules_list = []
    for module in MODULE_CONFIG.get('modules', []):
        modules_list.append({
            'id': module['id'],
            'name': module['name'],
            'description': module.get('description', ''),
            'learning_objectives': module.get('learning_objectives', []),
            'topics': module.get('relevant_topics', [])
        })
    
    return jsonify({
        'total_modules': len(modules_list),
        'modules': modules_list
    })


@app.route('/api/summarize', methods=['POST'])
def summarize():
    """
    Summarize a lecture document with module-specific adaptation
    
    Request body:
    - file: PDF/PPTX/TXT file
    - module_id: Module code (e.g., "SE3070")
    - use_ml_model: true/false (use ML model or fallback)
    
    Returns:
    - summary: List of selected sentences
    - flashcards: Generated Q&A flashcards
    - learning_objectives: Module's LO1-LO5
    - statistics: Coverage info
    """
    try:
        # Get parameters — frontend sends 'module' or 'module_id'
        module_id = request.form.get('module_id') or request.form.get('module', 'IT2040')
        use_ml_model = request.form.get('use_ml_model', 'true').lower() == 'true'
        
        # Extract file
        file = request.files.get('file')
        if not file:
            return jsonify({'error': 'No file provided'}), 400
        
        # Extract text
        filename = file.filename.lower()
        content = file.read()
        
        if filename.endswith('.pdf'):
            text = extract_text_from_pdf(content)
        elif filename.endswith('.pptx'):
            text = extract_text_from_pptx(content)
        elif filename.endswith('.txt'):
            text = content.decode('utf-8', errors='ignore')
        else:
            return jsonify({'error': 'Unsupported file type. Use PDF, PPTX, or TXT'}), 400
        
        if not text or len(text.strip()) < 50:
            return jsonify({'error': 'Insufficient text extracted from file'}), 400
        
        # Clean text — remove PDF/slide artifacts
        text = re.sub(r'/g\d+', '', text)
        text = re.sub(r'\(cid:\d+\)', '', text)
        text = re.sub(r'§', ' ', text)
        # Remove module-code header lines (pipe AND dash separators)
        text = re.sub(r'[A-Z]{2}\d{3,4}\s*[-–]\s*[^\n]*', ' ', text)
        text = re.sub(r'[A-Z]{2}\d{3,4}\s*\|[^\n]*', ' ', text)
        text = re.sub(r'IE\d{4}\s*\|[^\n]*', ' ', text)
        # Remove standalone "Lecture NN ..." header lines
        text = re.sub(r'\bLecture\s+\d+[^\n.!?]{0,80}\n', ' ', text)
        # Remove other slide header/footer patterns
        text = re.sub(r'Module\s+Code\s*\|[^.!?\n]*', ' ', text)
        text = re.sub(r'Module\s+Name\s*\|[^.!?\n]*', ' ', text)
        text = re.sub(r'Lecture\s+Title\s*\|[^.!?\n]*', ' ', text)
        text = re.sub(r'Lecturer\s*\.?\s*$', ' ', text, flags=re.MULTILINE)
        text = re.sub(r'Shashika\s*Lokuliyana\s*\d*', ' ', text, flags=re.IGNORECASE)
        text = re.sub(r'Faculty\s*of\s*Computing[^.]*Department\s+of\s+\w+', ' ', text)
        # Fix camelCase boundaries
        text = re.sub(r'(?<=[a-z])(?=[A-Z])', ' ', text)
        text = re.sub(r'\s+', ' ', text).strip()
        
        # Extract sentences
        sentences = extract_sentences_from_text(text)
        if len(sentences) < 3:
            return jsonify({'error': 'Could not extract enough sentences from text'}), 400
        
        # Get module config
        module_config = None
        if MODULE_CONFIG:
            for m in MODULE_CONFIG.get('modules', []):
                if m['id'] == module_id:
                    module_config = m
                    break
        
        if not module_config:
            return jsonify({'error': f'Module {module_id} not found in configuration'}), 400
        
        learning_objectives = module_config.get('learning_objectives', [])
        
        # Score sentences
        keyword_scores = np.array([
            sum(1 for kw in module_config.get('keywords', []) if kw.lower() in s.lower())
            for s in sentences
        ], dtype=float)
        if keyword_scores.max() > 0:
            keyword_scores = keyword_scores / keyword_scores.max()

        # Position-based scoring: earlier sentences in a lecture are often more important
        position_scores = np.array([
            1.0 - (0.3 * i / max(len(sentences) - 1, 1)) for i in range(len(sentences))
        ])

        if use_ml_model and MODEL is not None:
            ml_scores = score_sentences_with_model(sentences)
            if BOOSTER_AVAILABLE and MODULE_CONFIG:
                ml_scores = boost_scores_for_module(ml_scores, sentences, module_id, MODULE_CONFIG)
            scores = 0.65 * ml_scores + 0.25 * keyword_scores + 0.10 * position_scores
        else:
            scores = 0.6 * keyword_scores + 0.4 * position_scores
        
        # Select top sentences — aim for ~15-30 key sentences
        top_percentage = min(0.15, 25 / max(len(sentences), 1))
        n_select = max(5, min(30, int(len(sentences) * top_percentage)))
        top_indices = np.argsort(scores)[::-1][:n_select]
        top_indices = sorted(top_indices)  # Preserve original order
        
        selected_sentences = [sentences[i] for i in top_indices]
        selected_scores = scores[top_indices]
        
        # Create summary
        summary_text = " ".join(selected_sentences)
        
        # Generate flashcards
        flashcards = generate_flashcards_from_sentences(
            selected_sentences,
            module_id,
            learning_objectives
        )
        
        # Get statistics
        stats = {
            'total_sentences': len(sentences),
            'selected_sentences': len(selected_sentences),
            'coverage_percentage': round(len(selected_sentences) / len(sentences) * 100, 1),
            'average_score': round(float(np.mean(selected_scores)), 3),
            'score_range': [round(float(np.min(selected_scores)), 3), round(float(np.max(selected_scores)), 3)]
        }
        
        return jsonify({
            'module_id': module_id,
            'module_name': module_config.get('name', ''),
            'learning_objectives': learning_objectives,
            'summary': summary_text,
            'summary_sentences': selected_sentences,
            'flashcards': flashcards,
            'statistics': stats,
            'model_info': {
                'type': 'ML RandomForest + Module Adaptation',
                'ml_model_used': use_ml_model and MODEL is not None,
                'fallback_used': not (use_ml_model and MODEL is not None)
            }
        })
    
    except Exception as e:
        print(f"Error in summarize: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/api/flashcards', methods=['POST'])
def get_flashcards():
    """
    Generate flashcards from a summary
    
    Request body:
    - module_id: Module code
    - summary: Summary text or list of sentences
    
    Returns:
    - flashcards: List of question-answer pairs with LO tags
    """
    try:
        data = request.get_json() or {}
        module_id = data.get('module_id') or data.get('module', 'IT2040')
        summary = data.get('summary', '')
        
        if not summary:
            return jsonify({'error': 'No summary provided'}), 400
        
        # Get module config
        module_config = None
        if MODULE_CONFIG:
            for m in MODULE_CONFIG.get('modules', []):
                if m['id'] == module_id:
                    module_config = m
                    break
        
        if not module_config:
            return jsonify({'error': f'Module {module_id} not found'}), 400
        
        learning_objectives = module_config.get('learning_objectives', [])
        
        # Parse summary (could be string or list)
        if isinstance(summary, str):
            sentences = extract_sentences_from_text(summary)
        else:
            sentences = summary
        
        # Generate flashcards
        flashcards = generate_flashcards_from_sentences(
            sentences,
            module_id,
            learning_objectives
        )
        
        return jsonify({
            'module_id': module_id,
            'total_flashcards': len(flashcards),
            'flashcards': flashcards
        })
    
    except Exception as e:
        print(f"Error generating flashcards: {e}")
        return jsonify({'error': str(e)}), 500


# ─── Run ──────────────────────────────────────────────────────────────────────
if __name__ == '__main__':
    print("\n" + "=" * 70)
    print("🚀 AURA LEARN - Lecture Summarizer & Flashcard Generator")
    print("=" * 70)
    print(f"NLP: {'✅' if NLP_AVAILABLE else '❌'}")
    print(f"PDF: {'✅' if PDF_AVAILABLE else '❌'}")
    print(f"PPTX: {'✅' if PPTX_AVAILABLE else '❌'}")
    print(f"ML Model: {'✅' if MODEL is not None else '❌'}")
    print(f"Module Config: {'✅' if MODULE_CONFIG is not None else '❌'}")
    print("=" * 70)
    print("Starting server on http://localhost:5003")
    print("=" * 70 + "\n")
    
    app.run(host='0.0.0.0', port=5003, debug=True)
